#!/usr/bin/env python3
"""Read-only Worship source analyzer.

This script inspects one bulletin HWP plus one or more PPTX files and emits a
JSON candidate report. It does not write to the Mindex database.

The report is intentionally conservative:
- the bulletin/order sheet is treated as the service structure authority;
- PPTX files are treated as slide material;
- praise titles derived only from cropped PPT lyrics are marked as weak hints.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import zlib
from pathlib import Path
from typing import Any


TAG_PARA_TEXT = 67
CIRCLED_SERVICE_MARKERS = {"1": "①", "2": "②"}
ORDER_LABELS = {
  "환영",
  "경배와찬양",
  "경배와 찬양",
  "예배의부름",
  "예배의 부름",
  "사도신경",
  "신앙고백",
  "찬양",
  "찬송",
  "참회기도",
  "사죄의선언",
  "기도",
  "대표기도",
  "성경봉독",
  "말씀선포",
  "말씀",
  "설교",
  "결단의기도",
  "결단기도",
  "봉헌찬송",
  "봉헌기도",
  "교회소식",
  "새가족환영",
  "공동체고백",
  "파송찬송",
  "송영",
  "축도",
  "아멘송",
  "묵도",
}
PRAISE_LABELS = {
  "경배와찬양",
  "경배와 찬양",
  "찬양",
  "찬송",
  "봉헌찬송",
  "파송찬송",
  "송영",
  "아멘송",
}
NON_PRESENTED_ORDER_LABELS = {
  "환영",
  "사죄의선언",
  "아멘송",
}


def normalize_space(value: str) -> str:
  return re.sub(r"\s+", " ", value or "").strip()


def normalize_order_label(value: str) -> str:
  value = normalize_space(value)
  value = re.sub(r"^[※♱❦✽\s]+", "", value)
  value = value.replace(" ", "")
  return value


def canonical_order_label(value: str) -> str:
  normalized = normalize_order_label(value)
  aliases = {
    normalize_order_label("경배와찬양"): "찬양",
    normalize_order_label("경배와 찬양"): "찬양",
    normalize_order_label("말씀선포"): "설교",
    normalize_order_label("말씀"): "설교",
    normalize_order_label("결단의기도"): "결단기도",
  }
  return aliases.get(normalized, strip_presentation_markers(value))


def is_non_presented_order_label(value: str) -> bool:
  normalized = normalize_order_label(value)
  return normalized in {normalize_order_label(label) for label in NON_PRESENTED_ORDER_LABELS}


def split_order_parts(value: str) -> list[str]:
  return [normalize_space(part) for part in re.split(r"\s*/\s*", value) if normalize_space(part)]


def strip_presentation_markers(value: str) -> str:
  value = normalize_space(value)
  value = re.sub(r"^[※♱❦✽\s]+", "", value)
  return normalize_space(value)


def strip_quotes(value: str) -> str:
  value = normalize_space(value)
  return value.strip(" \"'“”‘’「」")


def strip_page_hint(value: str) -> str:
  value = re.sub(r"\([^)]*p\.[^)]+\)", "", value)
  return normalize_space(value)


def is_order_label(value: str) -> bool:
  return normalize_order_label(value) in {normalize_order_label(label) for label in ORDER_LABELS}


def is_praise_label(value: str) -> bool:
  return normalize_order_label(value) in {normalize_order_label(label) for label in PRAISE_LABELS}


def is_ignorable_order_text(value: str) -> bool:
  normalized = normalize_order_label(value)
  return normalized in {
    "인도자",
    "다같이",
    "말씀과봉헌",
    "파송과축복",
  }


def hangul_score(value: str) -> int:
  return len(re.findall(r"[가-힣]", value or ""))


def clean_hwp_text(value: str) -> str:
  value = value.replace("\x00", "")
  value = re.sub(r"[\x01-\x08\x0b\x0c\x0e-\x1f]", "", value)
  value = normalize_space(value)
  if not value:
    return ""
  # HWP records sometimes include binary-looking table control text such as
  # "氠瑢". Keep lines with real Korean content and drop obvious mojibake.
  if hangul_score(value) == 0 and not re.search(r"\d", value):
    return ""
  return value


def extract_hwp_texts(path: Path) -> list[str]:
  try:
    import olefile  # type: ignore
  except ImportError as exc:
    raise RuntimeError("olefile is required for HWP extraction") from exc

  ole = olefile.OleFileIO(str(path))
  flags = int.from_bytes(ole.openstream("FileHeader").read()[36:40], "little")
  compressed = bool(flags & 1)
  texts: list[str] = []

  section_names = ["/".join(parts) for parts in ole.listdir() if parts[:1] == ["BodyText"]]
  for name in sorted(section_names):
    raw = ole.openstream(name).read()
    if compressed:
      raw = zlib.decompress(raw, -15)
    pos = 0
    while pos + 4 <= len(raw):
      header = int.from_bytes(raw[pos:pos + 4], "little")
      pos += 4
      tag = header & 0x3ff
      size = (header >> 20) & 0xfff
      if size == 0xfff:
        size = int.from_bytes(raw[pos:pos + 4], "little")
        pos += 4
      data = raw[pos:pos + size]
      pos += size
      if tag != TAG_PARA_TEXT:
        continue
      text = clean_hwp_text(data.decode("utf-16le", errors="ignore"))
      if text:
        texts.append(text)
  return texts


def find_index(lines: list[str], pattern: str, start: int = 0) -> int | None:
  regex = re.compile(pattern)
  for index in range(start, len(lines)):
    if regex.search(lines[index]):
      return index
  return None


def slice_between(lines: list[str], start_pattern: str, end_pattern: str | None = None, start: int = 0) -> list[str]:
  begin = find_index(lines, start_pattern, start)
  if begin is None:
    return []
  if end_pattern is None:
    return lines[begin:]
  end = find_index(lines, end_pattern, begin + 1)
  return lines[begin:end if end is not None else len(lines)]


def make_element_from_order_line(line: str, service_key: str = "") -> dict[str, Any] | None:
  line = strip_presentation_markers(line)
  if not line:
    return None
  if line.startswith("▸"):
    return {
      "section": "기도문",
      "element_type": "body",
      "label": "기도문",
      "title": "기도문",
      "person": "",
      "body": line,
      "review_status": "matched",
      "confidence": "medium",
    }
  if line.startswith("▪"):
    return {
      "section": "찬양",
      "element_type": "praise",
      "label": "찬양",
      "title": strip_quotes(line.lstrip("▪").strip()),
      "person": "",
      "review_status": "needs_manual_praise",
      "confidence": "medium",
      "notes": ["bulletin song title hint; confirm/register manually before linking Praise"],
    }

  parts = split_order_parts(line)
  if not parts:
    return None
  label = strip_presentation_markers(parts[0])
  normalized = normalize_order_label(label)
  canonical_label = canonical_order_label(label)
  values = parts[1:]
  raw_value = " / ".join(values)
  if is_non_presented_order_label(label):
    return None

  if is_praise_label(label) or normalized.startswith(normalize_order_label("경배와찬양")):
    first_value = strip_quotes(values[0]) if values else ""
    title = "" if first_value in {"다같이", "다 같 이"} else first_value
    person = values[-1] if len(values) > 1 else ""
    return {
      "section": canonical_label,
      "element_type": "praise",
      "label": canonical_label,
      "title": title or canonical_label,
      "person": person if person not in {"다같이", "다 같 이"} else "",
      "review_status": "needs_manual_praise",
      "confidence": "medium" if title else "low",
      "notes": ["praise element from bulletin; user will confirm/register song directly"],
    }

  if normalized in {normalize_order_label(x) for x in ["예배의부름", "예배의 부름", "사도신경", "신앙고백"]}:
    title = "사도신경" if "사도신경" in line else canonical_label
    return {
      "section": "신앙고백" if "사도신경" in line else canonical_label,
      "element_type": "body",
      "label": canonical_label,
      "title": title,
      "person": values[-1] if values and values[-1] not in {"다같이", "다 같 이"} else "",
      "body": "",
      "review_status": "matched",
      "confidence": "high",
    }

  if normalized in {normalize_order_label(x) for x in ["기도", "대표기도", "참회기도", "결단의기도", "결단기도", "봉헌기도", "묵도"]}:
    person = values[-1] if values else ""
    return {
      "section": canonical_label,
      "element_type": "title_person",
      "label": canonical_label,
      "title": canonical_label,
      "person": "" if person in {"다같이", "다 같 이", "인도자", "인 도 자"} else person,
      "review_status": "matched" if values else "needs_review",
      "confidence": "high" if values else "medium",
    }

  if normalized == normalize_order_label("성경봉독"):
    reference = strip_page_hint(values[0]) if values else ""
    person = values[-1] if len(values) > 1 else ""
    return {
      "section": "성경봉독",
      "element_type": "scripture_reading",
      "label": "성경봉독",
      "title": "성경봉독",
      "person": "" if person in {"인도자", "인 도 자"} else person,
      "scripture_reference": reference,
      "review_status": "matched" if reference else "needs_review",
      "confidence": "high" if reference else "medium",
    }

  if normalized in {normalize_order_label("말씀선포"), normalize_order_label("설교"), normalize_order_label("말씀")}:
    title = ""
    person = ""
    for value in values:
      if re.search(r"[“”‘’「」]", value):
        title = strip_quotes(value)
      elif "목사" in value or "전도사" in value:
        person = value
    if not title and values:
      title = canonical_label if person else strip_quotes(values[0])
    return {
      "section": "설교",
      "element_type": "title_person",
      "label": canonical_label,
      "title": title or canonical_label,
      "person": person,
      "review_status": "matched" if title or person else "needs_review",
      "confidence": "high" if title and person else "medium",
    }

  if normalized in {normalize_order_label(x) for x in ["교회소식", "교회소식&새가족환영", "새가족환영", "공동체고백", "환영", "사죄의선언"]}:
    person = values[-1] if values else ""
    return {
      "section": canonical_label,
      "element_type": "title_person",
      "label": canonical_label,
      "title": canonical_label,
      "person": "" if person in {"다같이", "다 같 이", "인도자", "인 도 자"} else person,
      "review_status": "matched",
      "confidence": "high",
    }

  if normalized == normalize_order_label("축도"):
    person = ""
    notes = []
    for value in values:
      if value in {"인도자", "인 도 자"}:
        continue
      if "※" in value or "주기도문" in value:
        notes.append(value)
        continue
      person = value
    element = {
      "section": "축도",
      "element_type": "title_person",
      "label": "축도",
      "title": "축도",
      "person": person,
      "review_status": "matched" if values else "needs_review",
      "confidence": "high" if values else "medium",
    }
    if notes:
      element["notes"] = notes
    return element

  return {
    "section": label,
    "element_type": "editable",
    "label": label,
    "title": strip_quotes(raw_value) if raw_value else label,
    "person": "",
    "review_status": "needs_review",
    "confidence": "low",
    "notes": ["unclassified bulletin order line"],
  }


def merge_vertical_order_lines(lines: list[str]) -> list[str]:
  merged: list[str] = []
  index = 0
  while index < len(lines):
    line = strip_presentation_markers(lines[index])
    if not line:
      index += 1
      continue
    if line.startswith("▪"):
      merged.append(line)
      index += 1
      continue
    if is_order_label(line):
      values: list[str] = []
      lookahead = index + 1
      while lookahead < len(lines):
        candidate = strip_presentation_markers(lines[lookahead])
        if not candidate:
          lookahead += 1
          continue
        if candidate.startswith("▪") or is_order_label(candidate):
          break
        values.append(candidate)
        lookahead += 1
        # Most vertical bulletin cells use exactly one value line. Sermon rows
        # can still include title/person joined by slash in that one value.
        break
      merged.append(f"{line}/ {' / '.join(values)}" if values else line)
      index = lookahead
      continue
    merged.append(line)
    index += 1
  return merged


def filter_orderish_lines(lines: list[str]) -> list[str]:
  ignored_patterns = [
    r"^\d+부[‧/]",
    r"^3부/오전",
    r"^오전",
    r"^오후\s*\d",
    r"^\[다음주",
    r"^다음주",
    r"^※ 표시는",
    r"^표시는",
    r"^제\d+권",
    r"^이 름:",
    r"^주.?일.?오.?후.?예.?배$",
    r"^제자헌신예배$",
  ]
  filtered: list[str] = []
  for line in lines:
    clean = strip_presentation_markers(line)
    if not clean:
      continue
    if is_ignorable_order_text(clean):
      continue
    if any(re.search(pattern, clean) for pattern in ignored_patterns):
      continue
    filtered.append(clean)
  return filtered


def split_combined_public_lines(lines: list[str], service_number: str) -> list[str]:
  marker = CIRCLED_SERVICE_MARKERS[service_number]
  other_markers = [value for key, value in CIRCLED_SERVICE_MARKERS.items() if key != service_number]
  result: list[str] = []
  carry_label = ""
  for line in filter_orderish_lines(lines):
    if "/" in line and not line.startswith("/"):
      first = split_order_parts(line)[0]
      if first:
        carry_label = first
    if any(other in line for other in other_markers) and marker not in line:
      continue
    if marker in line:
      line = line.replace(marker, "")
      if line.startswith("/") and carry_label:
        line = f"{carry_label}{line}"
    result.append(line)
  return result


def elements_from_lines(lines: list[str], service_key: str, vertical: bool = False) -> list[dict[str, Any]]:
  working = merge_vertical_order_lines(filter_orderish_lines(lines)) if vertical else filter_orderish_lines(lines)
  elements: list[dict[str, Any]] = []
  for line in working:
    element = make_element_from_order_line(line, service_key)
    if element:
      element["source_line"] = line
      elements.append(element)
  return elements


def grouped_sections(elements: list[dict[str, Any]]) -> list[dict[str, Any]]:
  sections: list[dict[str, Any]] = []
  for element in elements:
    title = element.get("section") or element.get("label") or "기타"
    if not sections or sections[-1]["title"] != title:
      sections.append({"title": title, "elements": []})
    sections[-1]["elements"].append(element)
  return sections


def make_service_candidate(key: str, title: str, lines: list[str], vertical: bool = False) -> dict[str, Any]:
  elements = elements_from_lines(lines, key, vertical=vertical)
  manual_praise = [
    {
      "title_hint": element.get("title", ""),
      "label": element.get("label", ""),
      "source_line": element.get("source_line", ""),
      "reason": "찬양은 사용자가 직접 확인/등록 후 Mindex Praise에 연결",
    }
    for element in elements
    if element.get("review_status") == "needs_manual_praise"
  ]
  return {
    "service_key": key,
    "title": title,
    "source": "bulletin",
    "confidence": "draft",
    "sections": grouped_sections(elements),
    "manual_required": manual_praise,
    "stats": {
      "elements": len(elements),
      "manual_praise": len(manual_praise),
      "auto_extractable": len(elements) - len(manual_praise),
    },
  }


def build_bulletin_candidates(services: dict[str, list[str]]) -> list[dict[str, Any]]:
  afternoon_lines = services.get("sunday_afternoon", [])
  afternoon_title = "주일오후예배"
  if len(afternoon_lines) > 1 and "예배" in afternoon_lines[1]:
    afternoon_title = f"주일오후예배 - {strip_presentation_markers(afternoon_lines[1])}"
  return [
    make_service_candidate(
      "sunday_1st",
      "주일예배 (1부)",
      split_combined_public_lines(services.get("sunday_1_2_combined", []), "1"),
    ),
    make_service_candidate(
      "sunday_2nd",
      "주일예배 (2부)",
      split_combined_public_lines(services.get("sunday_1_2_combined", []), "2"),
    ),
    make_service_candidate(
      "sunday_3rd",
      "주일예배 (3부)",
      services.get("sunday_3rd", []),
      vertical=True,
    ),
    make_service_candidate(
      "sunday_afternoon",
      afternoon_title,
      afternoon_lines,
      vertical=True,
    ),
  ]


def parse_bulletin(path: Path) -> dict[str, Any]:
  lines = extract_hwp_texts(path)
  public_start = find_index(lines, r"3부/오전|1부.?오전")
  public_end = find_index(lines, r"\[다음주기도\]|이 름:", public_start or 0)
  public_lines = lines[public_start:public_end] if public_start is not None else []

  third_start = None
  if public_lines:
    praise_hits = [i for i, line in enumerate(public_lines) if re.search(r"경배와찬양|경배와 찬양", line)]
    if len(praise_hits) >= 2:
      third_start = praise_hits[1]
  third_lines = public_lines[third_start:] if third_start is not None else []
  first_second_lines = public_lines[:third_start] if third_start is not None else public_lines

  afternoon_lines = slice_between(lines, r"^주.*일.*오.*후.*예.*배$", r"✽다음주|위임목사", 0)
  friday_lines = slice_between(lines, r"^3금.*요.*기.*도.*회$|^“월삭예배”$", r"^월-금|^주.*일.*오.*후", 0)

  services = {
    "sunday_1_2_combined": first_second_lines,
    "sunday_3rd": third_lines,
    "sunday_afternoon": afternoon_lines,
    "friday_or_monthly_notice": friday_lines,
  }

  return {
    "path": str(path),
    "line_count": len(lines),
    "services": services,
    "service_candidates": build_bulletin_candidates(services),
    "raw_excerpt": [{"index": i + 1, "text": line} for i, line in enumerate(lines[:220])],
  }


def extract_pptx_slides(path: Path) -> list[dict[str, Any]]:
  try:
    from pptx import Presentation  # type: ignore
  except ImportError as exc:
    raise RuntimeError("python-pptx is required for PPTX extraction") from exc

  prs = Presentation(str(path))
  slides: list[dict[str, Any]] = []
  for index, slide in enumerate(list(prs.slides), 1):
    text_runs: list[str] = []
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        text = normalize_space(getattr(shape, "text", ""))
        if text:
          text_runs.append(text)
    text = " | ".join(text_runs)
    slides.append({
      "number": index,
      "text": text,
      "shape_text_count": len(text_runs),
      "class": classify_slide_text(text),
    })
  return slides


def classify_slide_text(text: str) -> str:
  text = normalize_space(text)
  if not text:
    return "blank_or_media"
  if "사도신경" in text or "전능하신 아버지" in text:
    return "faith_confession"
  if "성경봉독" in text:
    return "scripture_title"
  if re.match(r"^\d+\s", text) and re.search(r"(장|복음|에베소서|에스더|요한복음)", text):
    return "scripture_body"
  if "말씀선포" in text or "목사" in text and re.search(r"[‘「].+[’」]", text):
    return "sermon_title"
  if "봉헌기도" in text:
    return "offering_prayer"
  if "봉헌" in text:
    return "offering"
  if "교회소식" in text or "새가족" in text:
    return "announcement"
  if "축도" in text:
    return "benediction"
  if text.startswith("♪"):
    return "praise_title"
  if re.match(r"^\d+[.)]\s*", text) or re.search(r"\|\s*\d+[.)]", text):
    return "lyrics_or_hymn_body"
  return "body_or_title"


def block_key(slide: dict[str, Any]) -> str:
  text = slide["text"]
  if slide["class"] == "blank_or_media":
    return "blank_or_media"
  # Repeated hymn slides often have identical multi-verse text. Collapse those.
  return f"{slide['class']}:{normalize_space(text)[:80]}"


def group_slide_blocks(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
  blocks: list[dict[str, Any]] = []
  current: dict[str, Any] | None = None
  for slide in slides:
    key = block_key(slide)
    if current and current["key"] == key:
      current["end"] = slide["number"]
      current["count"] += 1
      continue
    if current:
      blocks.append(current)
    current = {
      "key": key,
      "start": slide["number"],
      "end": slide["number"],
      "count": 1,
      "class": slide["class"],
      "sample": slide["text"][:220],
    }
  if current:
    blocks.append(current)
  return blocks


def infer_service_hint(path: Path, slides: list[dict[str, Any]], bulletin: dict[str, Any] | None) -> str:
  name = path.stem.lower()
  text_blob = "\n".join(slide["text"] for slide in slides[:120])
  if "아름다운 이름" in text_blob or "엡 1:1" in text_blob:
    return "sunday_afternoon"
  if "곤한 내 영혼" in text_blob or "할렐루야 찬양대" in text_blob:
    return "sunday_3rd"
  if "요 19:38" in text_blob or "_1st" in name:
    return "sunday_1st"
  if "에 9:20" in text_blob or "_2nd" in name:
    return "sunday_2nd"
  return "needs_review"


def summarize_pptx(path: Path, bulletin: dict[str, Any] | None = None) -> dict[str, Any]:
  slides = extract_pptx_slides(path)
  blocks = group_slide_blocks(slides)
  class_counts: dict[str, int] = {}
  for slide in slides:
    class_counts[slide["class"]] = class_counts.get(slide["class"], 0) + 1
  return {
    "path": str(path),
    "slide_count": len(slides),
    "service_hint": infer_service_hint(path, slides, bulletin),
    "class_counts": class_counts,
    "blocks": blocks,
    "first_slides": slides[:16],
    "important_slides": [
      slide for slide in slides
      if slide["class"] not in {"blank_or_media", "lyrics_or_hymn_body", "body_or_title"}
    ][:80],
    "notes": [
      "PPT praise titles are hints only. Cropped first-line lyric titles must be matched against Mindex Praise or a separate setlist.",
      "Filename suffix is not authoritative; service_hint uses content fingerprints.",
    ],
  }


def main() -> int:
  parser = argparse.ArgumentParser(description="Analyze Worship bulletin/PPTX sources without writing DB data.")
  parser.add_argument("--bulletin-hwp", type=Path)
  parser.add_argument("--pptx", type=Path, action="append", default=[])
  parser.add_argument("--pretty", action="store_true")
  args = parser.parse_args()

  bulletin = parse_bulletin(args.bulletin_hwp) if args.bulletin_hwp else None
  report = {
    "bulletin": bulletin,
    "pptx": [summarize_pptx(path, bulletin) for path in args.pptx],
    "algorithm": {
      "authority_order": ["bulletin_structure", "manual_praise_confirmation", "mindex_praise", "mindex_scripture", "pptx_slide_material"],
      "candidate_flow": [
        "extract bulletin cover/next-page order",
        "split service candidates by order markers",
        "extract non-praise worship elements from bulletin into service candidates",
        "mark praise elements as manual confirmation/registration targets",
        "extract PPT slide text and group repeated slide blocks",
        "infer service by content fingerprints, not only filename",
        "never treat cropped PPT praise titles as canonical",
        "emit import candidates for review before writing canonical Worship tables",
      ],
    },
  }
  json.dump(report, sys.stdout, ensure_ascii=False, indent=2 if args.pretty else None)
  sys.stdout.write("\n")
  return 0


if __name__ == "__main__":
  raise SystemExit(main())
